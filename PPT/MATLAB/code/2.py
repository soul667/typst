import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_ppt(pdf_path, ppt_path):
    # 打开PDF文件
    doc = fitz.open(pdf_path)
    
    # 创建一个PPT文件
    pres = Presentation()
    
    # 遍历PDF中的每一页
    for page_number in range(doc.page_count):
        page = doc[page_number]
        
        # 将页面渲染为图像
        pix = page.get_pixmap()
        img_path = f"page_{page_number}.png"
        pix.save(img_path)
        
        # 创建一个新的幻灯片
        slide = pres.slides.add_slide(pres.slide_layouts[5])
        
        # 在幻灯片中添加图像
        img = slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8))
        
        # 删除临时保存的图像文件
        os.remove(img_path)
    
    # 保存PPT文件
    pres.save(ppt_path)

# 指定PDF和PPT的路径
pdf_path = './研讨.pdf'
ppt_path = './output.pptx'

# 调用函数
pdf_to_ppt(pdf_path, ppt_path)
